
import os
import json
import zipfile
import pandas as pd
import PyPDF2
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
from pptx import Presentation
from docx import Document
import io
from typing import Dict, Any, Optional, Tuple, List
import asyncio
import xml.etree.ElementTree as ET
import re


class FileProcessor:


                       
    SUPPORTED_EXTENSIONS = {
                    
        '.txt', '.py', '.md', '.rst', '.log', '.cfg', '.ini', '.yaml', '.yml',
                     
        '.xlsx', '.xls', '.csv', '.tsv',
                        
        '.docx', '.pdf', '.pptx',
                    
        '.json', '.jsonl', '.jsonld', '.pdb', '.xml', '.html', '.htm',
                          
        '.zip',
    }

    def __init__(self, base_dir: str = ""):

        self.base_dir = base_dir
        self.cache_dir = ""
        self.cache_file = ""
        self._processed_cache: Dict[str, str] = {}

               
        self._processors = {
            '.txt': self._process_text_file,
            '.py': self._process_text_file,
            '.md': self._process_text_file,
            '.rst': self._process_text_file,
            '.log': self._process_text_file,
            '.cfg': self._process_text_file,
            '.ini': self._process_text_file,
            '.yaml': self._process_text_file,
            '.yml': self._process_text_file,
            '.xlsx': self._process_excel_file,
            '.xls': self._process_excel_file,
            '.csv': self._process_csv_file,
            '.tsv': self._process_tsv_file,
            '.docx': self._process_docx_file,
            '.pdf': self._process_pdf_file,
            '.pptx': self._process_pptx_file,
            '.json': self._process_json_file,
            '.jsonl': self._process_jsonl_file,
            '.jsonld': self._process_jsonld_file,
            '.pdb': self._process_pdb_file,
            '.xml': self._process_xml_file,
            '.html': self._process_html_file,
            '.htm': self._process_html_file,
            '.zip': self._process_zip_file,
        }

    def check_cache(self, file_name: str) -> Optional[str]:

        return self._processed_cache.get(file_name)

    def _save_cache(self) -> None:

        if not self.cache_file:
            return
        try:
            with open(self.cache_file, "w", encoding="utf-8") as f:
                json.dump(self._processed_cache, f, ensure_ascii=False, separators=(',', ':'))
        except Exception:
            pass

    def set_base_dir(self, new_base_dir: str) -> None:

        if not new_base_dir:
            return
        self.base_dir = new_base_dir
        self.cache_dir = os.path.join(self.base_dir, "cache")
        self.cache_file = os.path.join(self.cache_dir, "processed_files.json")
        os.makedirs(self.cache_dir, exist_ok=True)
                             
        self._load_cache()

    def _load_cache(self) -> None:

        if not self.cache_file or not os.path.exists(self.cache_file):
            return
        try:
            with open(self.cache_file, "r", encoding="utf-8") as f:
                on_disk = json.load(f)
                if isinstance(on_disk, dict):
                                  
                    self._processed_cache = {**on_disk, **self._processed_cache}
        except Exception:
            pass

    def _process_text_file(self, file_name: str) -> str:

        file_path = os.path.join(self.base_dir, file_name)
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'gbk', 'gb2312']

        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                return f"=== Text File: {file_name} ===\n{content}"
            except UnicodeDecodeError:
                continue
            except Exception as e:
                return f"Error reading text file: {str(e)}"

                   
        try:
            with open(file_path, 'rb') as f:
                content = f.read().decode('utf-8', errors='replace')
            return f"=== Text File: {file_name} ===\n{content}"
        except Exception as e:
            return f"Error reading text file: {str(e)}"

    def _process_excel_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            excel_file = pd.ExcelFile(file_path)
            result = [f"=== Excel File: {file_name} ==="]

            for sheet_name in excel_file.sheet_names:
                              
                df = pd.read_excel(file_path, sheet_name=sheet_name, header=None)

                result.append(f"\n--- Sheet: {sheet_name} ---")
                result.append(f"Shape: {df.shape[0]} rows × {df.shape[1]} columns")

                      
                if df.shape[0] > 0:
                    first_row = df.iloc[0].astype(str).tolist()
                    if all(not str(v).replace('.','').replace('-','').isdigit() for v in first_row if pd.notna(v) and str(v).strip()):
                        df.columns = df.iloc[0]
                        df = df.iloc[1:]
                        result.append(f"Columns: {list(df.columns)}")

                                    
                result.append(f"\nData:\n{df.to_string(index=True, max_rows=None, max_cols=None)}")

                       
                numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
                if numeric_cols:
                    result.append(f"\n--- Statistics ---")
                    result.append(df[numeric_cols].describe().to_string())

            return "\n".join(result)
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"

    def _process_csv_file(self, file_name: str) -> str:

        file_path = os.path.join(self.base_dir, file_name)

                     
        for encoding in ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']:
            for sep in [',', ';', '\t', '|']:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, sep=sep)
                    if df.shape[1] > 1:            
                        break
                except:
                    continue
            else:
                continue
            break
        else:
                      
            try:
                df = pd.read_csv(file_path)
            except Exception as e:
                return f"Error reading CSV file: {str(e)}"

        result = [
            f"=== CSV File: {file_name} ===",
            f"Shape: {df.shape[0]} rows × {df.shape[1]} columns",
            f"Columns: {list(df.columns)}",
            f"\nData:\n{df.to_string(index=True, max_rows=None, max_cols=None)}"
        ]

              
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        if numeric_cols:
            result.append(f"\n--- Statistics ---")
            result.append(df[numeric_cols].describe().to_string())

        return "\n".join(result)

    def _process_tsv_file(self, file_name: str) -> str:

        file_path = os.path.join(self.base_dir, file_name)
        try:
            df = pd.read_csv(file_path, sep='\t')
            result = [
                f"=== TSV File: {file_name} ===",
                f"Shape: {df.shape[0]} rows × {df.shape[1]} columns",
                f"Columns: {list(df.columns)}",
                f"\nData:\n{df.to_string(index=True, max_rows=None, max_cols=None)}"
            ]
            return "\n".join(result)
        except Exception as e:
            return f"Error reading TSV file: {str(e)}"
    
    def _process_docx_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            doc = Document(file_path)
            result = [f"=== Word Document: {file_name} ==="]

                    
            for i, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    result.append(text)

                    
            for i, table in enumerate(doc.tables):
                result.append(f"\n--- Table {i+1} ---")
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(" | ".join(row_data))
                result.extend(table_data)

            return "\n".join(result)
        except Exception as e:
            return f"Error reading Word document: {str(e)}"

    def _process_pdf_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            text_content = []
            
            # Try pdfplumber first (better for tables and structured content)
            if PDFPLUMBER_AVAILABLE:
                try:
                    with pdfplumber.open(file_path) as pdf:
                        total_pages = len(pdf.pages)
                        text_content.append(f"=== PDF File: {file_name} ({total_pages} pages) ===")
                        
                        for page_num, page in enumerate(pdf.pages, 1):
                            try:
                                # Extract text
                                page_text = page.extract_text()
                                
                                # Extract tables
                                tables = page.extract_tables()
                                
                                page_content = []
                                if page_text and page_text.strip():
                                    page_content.append(page_text.strip())
                                
                                # Add tables in markdown format
                                if tables:
                                    for table_idx, table in enumerate(tables, 1):
                                        if table and len(table) > 0:
                                            page_content.append(f"\n--- Table {table_idx} on Page {page_num} ---")
                                            # Convert table to markdown
                                            if len(table) > 0:
                                                # Header row
                                                headers = table[0]
                                                if headers:
                                                    page_content.append("| " + " | ".join(str(cell) if cell else "" for cell in headers) + " |")
                                                    page_content.append("|" + "|".join(["---"] * len(headers)) + "|")
                                                # Data rows
                                                for row in table[1:]:
                                                    if row:
                                                        page_content.append("| " + " | ".join(str(cell) if cell else "" for cell in row) + " |")
                                
                                if page_content:
                                    text_content.append(f"\n--- Page {page_num} ---")
                                    text_content.append("\n".join(page_content))
                            except Exception:
                                continue
                        
                        full_text = "\n".join(text_content)
                        if len(full_text) > 200:
                            return full_text
                        # If pdfplumber didn't extract much, fall back to PyPDF2
                except Exception:
                    pass  # Fall back to PyPDF2
            
            # Fallback to PyPDF2
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                total_pages = len(pdf_reader.pages)
                text_content.append(f"=== PDF File: {file_name} ({total_pages} pages) ===")

                                     
                for page_num in range(total_pages):
                    try:
                        page = pdf_reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_content.append(f"\n--- Page {page_num + 1} ---")
                            text_content.append(page_text.strip())
                    except Exception:
                        continue

            full_text = "\n".join(text_content)

                              
            if len(full_text) < 200 and total_pages > 0:
                full_text += "\n\n[Note: PDF may contain scanned images. Text extraction may be incomplete. Consider using OCR if needed.]"

            return full_text
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"

    def _process_pptx_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            prs = Presentation(file_path)
            result = [f"=== PowerPoint: {file_name} ({len(prs.slides)} slides) ==="]

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    result.append(f"\n--- Slide {slide_num} ---")
                    result.extend(slide_text)

            return "\n".join(result)
        except Exception as e:
            return f"Error reading PowerPoint file: {str(e)}"

    def _process_json_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return f"=== JSON File: {file_name} ===\n{json.dumps(data, indent=2, ensure_ascii=False)}"
        except Exception as e:
            return f"Error reading JSON file: {str(e)}"

    def _process_jsonl_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            result = [f"=== JSONL File: {file_name} ==="]
            with open(file_path, 'r', encoding='utf-8') as f:
                for i, line in enumerate(f, 1):
                    line = line.strip()
                    if line:
                        try:
                            data = json.loads(line)
                            result.append(f"Line {i}: {json.dumps(data, ensure_ascii=False)}")
                        except json.JSONDecodeError:
                            result.append(f"Line {i}: {line}")
            return "\n".join(result)
        except Exception as e:
            return f"Error reading JSONL file: {str(e)}"

    def _process_jsonld_file(self, file_name: str) -> str:

        return self._process_json_file(file_name)

    def _process_pdb_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return f"=== PDB File: {file_name} ===\n{content}"
        except Exception as e:
            return f"Error reading PDB file: {str(e)}"

    def _process_xml_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            tree = ET.parse(file_path)
            root = tree.getroot()

            result = [f"=== XML File: {file_name} ==="]
            result.append(f"Root element: {root.tag}")
            if root.attrib:
                result.append(f"Root attributes: {root.attrib}")

            def process_element(element, level=0):
                indent = "  " * level
                element_info = f"{indent}<{element.tag}"

                if element.attrib:
                    attrs_str = " ".join([f'{k}="{v}"' for k, v in element.attrib.items()])
                    element_info += f" {attrs_str}"

                if element.text and element.text.strip():
                    text = element.text.strip()
                    element_info += f">{text}</{element.tag}>"
                elif len(element) == 0:
                    element_info += "/>"
                else:
                    element_info += ">"

                result.append(element_info)

                for child in element:
                    process_element(child, level + 1)

                if len(element) > 0:
                    result.append(f"{indent}</{element.tag}>")

            process_element(root)
            return "\n".join(result)
        except ET.ParseError as e:
            return f"Error parsing XML file: {str(e)}"
        except Exception as e:
            return f"Error reading XML file: {str(e)}"

    def _process_html_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

                               
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'html.parser')
                                
                for element in soup.find_all(['script', 'style']):
                    element.decompose()
                text = soup.get_text(separator='\n', strip=True)
                return f"=== HTML File: {file_name} ===\n{text}"
            except ImportError:
                             
                clean_text = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL | re.IGNORECASE)
                clean_text = re.sub(r'<style[^>]*>.*?</style>', '', clean_text, flags=re.DOTALL | re.IGNORECASE)
                clean_text = re.sub(r'<[^>]+>', ' ', clean_text)
                clean_text = re.sub(r'\s+', ' ', clean_text).strip()
                return f"=== HTML File: {file_name} ===\n{clean_text}"
        except Exception as e:
            return f"Error reading HTML file: {str(e)}"

    def _process_zip_file(self, file_name: str) -> str:

        try:
            file_path = os.path.join(self.base_dir, file_name)
            result = [f"=== ZIP File: {file_name} ==="]

            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                file_list = zip_ref.namelist()
                result.append(f"Total files: {len(file_list)}")
                result.append("Files:")
                for name in file_list:
                    result.append(f"  - {name}")

                safe_root = os.path.abspath(self.base_dir)
                safe_members = []
                skipped = []
                for name in file_list:
                    try:
                        dest_path = os.path.abspath(os.path.join(self.base_dir, name))
                        if os.path.commonpath([safe_root, dest_path]) != safe_root:
                            skipped.append(name)
                            continue
                    except Exception:
                        skipped.append(name)
                        continue
                    safe_members.append(name)

                for member in safe_members:
                    zip_ref.extract(member, self.base_dir)

                result.append(f"\nExtracted {len(safe_members)} files to: {self.base_dir}")
                if skipped:
                    preview = "\n".join(f"  - {name}" for name in skipped[:10])
                    extra = f"\n  ... and {len(skipped) - 10} more" if len(skipped) > 10 else ""
                    result.append(f"\nSkipped {len(skipped)} unsafe entries:\n{preview}{extra}")

            return "\n".join(result)
        except Exception as e:
            return f"Error reading ZIP file: {str(e)}"

    def process_file(self, file_name: str) -> str:

                           
        cached = self.check_cache(file_name)
        if cached is not None:
            return cached

        file_path = os.path.join(self.base_dir, file_name)
        if not os.path.exists(file_path):
            return f"Error: File not found: {file_name}"

        file_extension = os.path.splitext(file_name)[1].lower()

        if file_extension in self._processors:
            try:
                converted = self._processors[file_extension](file_name)
                               
                self._processed_cache[file_name] = converted
                self._save_cache()
                return converted
            except Exception as e:
                return f"Error processing {file_extension} file: {str(e)}"
        else:
                        
            try:
                return self._process_text_file(file_name)
            except Exception:
                return f"Unsupported file type: {file_extension}\nSupported types: {', '.join(sorted(self._processors.keys()))}"

    def list_files(self) -> List[str]:

        if not os.path.exists(self.base_dir):
            return []

        files = []
        for root, dirs, filenames in os.walk(self.base_dir):
                                  
            dirs[:] = [d for d in dirs if d != 'cache']
            for fname in filenames:
                abs_path = os.path.join(root, fname)
                rel_path = os.path.relpath(abs_path, self.base_dir)
                files.append(rel_path)
        return files


                    
_global_processor: Optional[FileProcessor] = None


def get_file_processor(base_dir: str = "") -> FileProcessor:

    global _global_processor
    if _global_processor is None:
        _global_processor = FileProcessor(base_dir)
    elif base_dir and _global_processor.base_dir != base_dir:
        _global_processor.set_base_dir(base_dir)
    return _global_processor


async def process_file_content(file_name: str, base_dir: str = "./data/files") -> str:










    try:
        processor = get_file_processor(base_dir)
        processor.set_base_dir(base_dir)

                                            
        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(None, processor.process_file, file_name)
        return result
    except Exception as e:
        return f"Error processing file: {str(e)}"


async def list_files_in_directory(directory_path: str) -> str:

    try:
        if not os.path.exists(directory_path):
            return f"Error: Directory not found: {directory_path}"

        if not os.path.isdir(directory_path):
            return f"Error: Path is not a directory: {directory_path}"

        files = []
        for item in os.listdir(directory_path):
            item_path = os.path.join(directory_path, item)
            if os.path.isfile(item_path):
                file_size = os.path.getsize(item_path)
                file_ext = os.path.splitext(item)[1].lower()
                files.append(f"  {item} ({file_ext}) - {file_size} bytes")

        if not files:
            return f"Directory '{directory_path}' is empty"

        result = [f"Files in '{directory_path}':", f"Total files: {len(files)}"]
        result.extend(sorted(files))

        return "\n".join(result)
    except Exception as e:
        return f"Error listing directory: {str(e)}"


def get_openai_function_process_file():

    return {
        "type": "function",
        "function": {
            "name": "process_file",
            "description": "Process and extract content from various file formats including text files, spreadsheets, documents, data files, and compressed archives. This tool reads the COMPLETE file content to ensure no information is missed.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_name": {
                        "type": "string",
                        "description": "The file name to process (relative to data directory).",
                        "examples": ["example.xlsx", "report.pdf", "data.csv", "config.xml"]
                    }
                },
                "required": ["file_name"]
            }
        }
    }


if __name__ == '__main__':
    import sys
    if len(sys.argv) > 1:
        processor = FileProcessor()
        processor.set_base_dir("./data/files")
        result = processor.process_file(sys.argv[1])
        print(result)

